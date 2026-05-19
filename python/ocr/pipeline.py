"""
ocr/pipeline.py — Text extraction
Automatically finds Tesseract on Windows
"""
import os
import logging
from pathlib import Path

log = logging.getLogger(__name__)

# ── Auto-find Tesseract on Windows ───────────────────────────────────────────
TESSERACT_AVAILABLE = False
try:
    import pytesseract
    # Common Windows install paths
    TESS_PATHS = [
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        r'C:\Users\admin\AppData\Local\Programs\Tesseract-OCR\tesseract.exe',
        r'C:\tools\Tesseract-OCR\tesseract.exe',
    ]
    for p in TESS_PATHS:
        if os.path.exists(p):
            pytesseract.pytesseract.tesseract_cmd = p
            log.info(f"Tesseract found at: {p}")
            break
    pytesseract.get_tesseract_version()
    TESSERACT_AVAILABLE = True
    log.info("Tesseract ready")
except Exception:
    log.warning("Tesseract not found — scanned image OCR disabled. Text PDFs still work fine.")


def extract_text(file_path: str, mime_type: str) -> dict:
    path = Path(file_path)
    ext  = path.suffix.lower()
    if not path.exists():
        return _err(f"File not found: {file_path}")
    if mime_type == "application/pdf" or ext == ".pdf":
        return _extract_pdf(file_path)
    if mime_type.startswith("image/") or ext in {".jpg", ".jpeg", ".png", ".tiff", ".bmp"}:
        return _extract_image(file_path)
    if ext == ".pptx" or "presentationml" in mime_type:
        return _extract_pptx(file_path)
    if ext in {".txt", ".md"}:
        return {"text": path.read_text(encoding="utf-8", errors="ignore"), "pages": 1, "method": "plaintext", "warning": None}
    return _err(f"Unsupported: {ext}")


def _extract_pdf(path: str) -> dict:
    try:
        import pdfplumber
        texts, pages = [], 0
        with pdfplumber.open(path) as pdf:
            pages = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text(x_tolerance=2, y_tolerance=2)
                if t and t.strip():
                    texts.append(t)
                elif TESSERACT_AVAILABLE:
                    try:
                        img = page.to_image(resolution=200).original
                        texts.append(_tesseract(img))
                    except Exception as e:
                        log.warning(f"OCR page failed: {e}")
        full = "\n\n".join(texts)
        if not full.strip():
            return {"text": "", "pages": pages, "method": "pdfplumber",
                    "warning": "Scanned PDF — install Tesseract or paste text manually"}
        return {"text": full, "pages": pages, "method": "pdfplumber", "warning": None}
    except Exception as e:
        log.error(f"PDF failed: {e}", exc_info=True)
        return _err(str(e))


def _extract_image(path: str) -> dict:
    if not TESSERACT_AVAILABLE:
        return {"text": "", "pages": 1, "method": "none",
                "warning": "Tesseract not found — paste text manually"}
    try:
        from PIL import Image
        img = Image.open(path)
        w, h = img.size
        if w < 1500:
            img = img.resize((int(w * 1500/w), int(h * 1500/w)), Image.LANCZOS)
        text = _tesseract(img)
        return {"text": text, "pages": 1, "method": "tesseract",
                "warning": None if text.strip() else "No text detected"}
    except Exception as e:
        return _err(str(e))


def _extract_pptx(path: str) -> dict:
    try:
        from pptx import Presentation
        texts = []
        prs = Presentation(path)
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t: parts.append(t)
            if parts: texts.append("\n".join(parts))
        return {"text": "\n\n".join(texts), "pages": len(prs.slides), "method": "pptx", "warning": None}
    except Exception as e:
        return _err(str(e))


def _tesseract(img) -> str:
    import pytesseract
    return pytesseract.image_to_string(img, lang="eng", config="--oem 3 --psm 6")


def _err(msg: str) -> dict:
    return {"text": "", "pages": 0, "method": "error", "warning": msg}
